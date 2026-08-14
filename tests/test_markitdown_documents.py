"""Regression: MarkItDown document conversion (used by
handlers/meeting_dev_handler.py::_scan_documents) previously had zero real
coverage — only mocked in test_meeting_dev_handler.py. These exercise the
real `MarkItDown().convert()` call against fully synthetic, generated-at-test-time
fixtures (no real client documents, no binaries committed to the repo).

PDF is intentionally omitted: hand-crafting a minimal PDF that pdfminer-six
(MarkItDown's PDF backend) reliably extracts text from is fragile without a
PDF-writing library, and adding one (e.g. reportlab) only to cover this one
format isn't worth the new dependency for a synthetic-fixture test.
"""

import zipfile

import pytest
from markitdown import MarkItDown

SYNTHETIC_TEXT = "Synthetic meeting attachment / Project: Northwind / Task: TEST-123"


@pytest.fixture
def md() -> MarkItDown:
    return MarkItDown()


def test_converts_synthetic_xlsx(tmp_path, md):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws["A1"] = SYNTHETIC_TEXT
    path = tmp_path / "sample.xlsx"
    wb.save(path)

    result = md.convert(str(path))
    assert SYNTHETIC_TEXT in result.text_content


def test_converts_synthetic_pptx(tmp_path, md):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    textbox.text_frame.text = SYNTHETIC_TEXT
    path = tmp_path / "sample.pptx"
    prs.save(path)

    result = md.convert(str(path))
    assert SYNTHETIC_TEXT in result.text_content


def test_converts_synthetic_docx(tmp_path, md):
    """python-docx isn't a project dependency (markitdown uses `mammoth` to
    read .docx, not python-docx) — building the file by hand via `zipfile`
    avoids adding a new dependency just for this fixture. A .docx is just a
    small OOXML zip: content-types manifest, package relationships, and the
    document body."""
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    package_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="word/document.xml"/>'
        "</Relationships>"
    )
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body><w:p><w:r><w:t>{SYNTHETIC_TEXT}</w:t></w:r></w:p></w:body>"
        "</w:document>"
    )

    path = tmp_path / "sample.docx"
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", package_rels)
        zf.writestr("word/document.xml", document_xml)

    result = md.convert(str(path))
    assert SYNTHETIC_TEXT in result.text_content
